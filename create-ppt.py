#!/usr/bin/env python3
"""
ChatSVTR 功能介绍 PPT 生成脚本
使用 python-pptx 库创建专业演示文稿
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# 定义 SVTR 品牌色
SVTR_BLUE = RGBColor(41, 98, 255)  # 主色调
SVTR_DARK = RGBColor(26, 32, 44)   # 深色
SVTR_LIGHT = RGBColor(237, 242, 247)  # 浅色

# ========== 幻灯片 1: 封面 ==========
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

# 背景色
background = slide1.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = SVTR_BLUE

# 标题
title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "ChatSVTR"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(72)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(255, 255, 255)
title_para.alignment = PP_ALIGN.CENTER

# 副标题
subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.8))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "SVTR.AI 全球AI创投平台 - 产品功能介绍"
subtitle_para = subtitle_frame.paragraphs[0]
subtitle_para.font.size = Pt(28)
subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
subtitle_para.alignment = PP_ALIGN.CENTER

# 底部信息
footer_box = slide1.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
footer_frame = footer_box.text_frame
footer_frame.text = "硅谷科技评论 | 2025"
footer_para = footer_frame.paragraphs[0]
footer_para.font.size = Pt(16)
footer_para.font.color.rgb = RGBColor(200, 220, 255)
footer_para.alignment = PP_ALIGN.CENTER

# ========== 幻灯片 2: 项目概述 ==========
slide2 = prs.slides.add_slide(prs.slide_layouts[6])

# 标题
title_box2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
title_frame2 = title_box2.text_frame
title_frame2.text = "项目概述"
title_para2 = title_frame2.paragraphs[0]
title_para2.font.size = Pt(44)
title_para2.font.bold = True
title_para2.font.color.rgb = SVTR_DARK

# 内容框
content_box2 = slide2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(4.5))
text_frame2 = content_box2.text_frame
text_frame2.word_wrap = True

items = [
    ("定位", "全球AI创投行业的统一平台"),
    ("核心价值", "连接AI创业者、投资人与行业专家"),
    ("技术特点", "混合RAG架构 + Cloudflare边缘计算"),
    ("数据来源", "飞书知识库 252个节点完整内容"),
    ("性能优化", "37.9KB资源优化 | Lighthouse 90+分")
]

for label, desc in items:
    p = text_frame2.add_paragraph()
    p.text = f"• {label}: {desc}"
    p.font.size = Pt(20)
    p.font.color.rgb = SVTR_DARK
    p.space_after = Pt(12)

# ========== 幻灯片 3: 核心技术栈 ==========
slide3 = prs.slides.add_slide(prs.slide_layouts[6])

title_box3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
title_frame3 = title_box3.text_frame
title_frame3.text = "核心技术栈"
title_para3 = title_frame3.paragraphs[0]
title_para3.font.size = Pt(44)
title_para3.font.bold = True
title_para3.font.color.rgb = SVTR_DARK

# 左侧：前端技术
left_box = slide3.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(4), Inches(4))
left_frame = left_box.text_frame
left_frame.word_wrap = True

p_left_title = left_frame.paragraphs[0]
p_left_title.text = "前端技术"
p_left_title.font.size = Pt(24)
p_left_title.font.bold = True
p_left_title.font.color.rgb = SVTR_BLUE

frontend = ["原生 HTML5/CSS3/JavaScript", "ES2022 语法标准", "WebP 图片优化", "Gzip 压缩"]
for item in frontend:
    p = left_frame.add_paragraph()
    p.text = f"  • {item}"
    p.font.size = Pt(18)
    p.space_after = Pt(8)

# 右侧：后端技术
right_box = slide3.shapes.add_textbox(Inches(5.2), Inches(1.8), Inches(4), Inches(4))
right_frame = right_box.text_frame
right_frame.word_wrap = True

p_right_title = right_frame.paragraphs[0]
p_right_title.text = "后端技术"
p_right_title.font.size = Pt(24)
p_right_title.font.bold = True
p_right_title.font.color.rgb = SVTR_BLUE

backend = ["Cloudflare Workers", "Cloudflare KV 存储", "Cloudflare Vectorize", "OpenAI GPT + Workers AI"]
for item in backend:
    p = right_frame.add_paragraph()
    p.text = f"  • {item}"
    p.font.size = Pt(18)
    p.space_after = Pt(8)

# ========== 幻灯片 4: AI & RAG 系统 ==========
slide4 = prs.slides.add_slide(prs.slide_layouts[6])

title_box4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
title_frame4 = title_box4.text_frame
title_frame4.text = "AI & RAG 智能问答系统"
title_para4 = title_frame4.paragraphs[0]
title_para4.font.size = Pt(44)
title_para4.font.bold = True
title_para4.font.color.rgb = SVTR_DARK

content_box4 = slide4.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(5))
text_frame4 = content_box4.text_frame
text_frame4.word_wrap = True

rag_features = [
    ("混合RAG架构", "结合向量搜索与关键词检索"),
    ("数据源", "飞书知识库 252 节点完整内容"),
    ("AI模型", "OpenAI GPT + Cloudflare Workers AI"),
    ("主服务", "functions/lib/hybrid-rag-service.ts"),
    ("飞书集成", "App ID: cli_a8e2014cbe7d9013"),
    ("智能同步", "scripts/smart-sync-strategy.js")
]

for label, desc in rag_features:
    p = text_frame4.add_paragraph()
    p.text = f"🔹 {label}"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = SVTR_BLUE
    p.space_after = Pt(6)

    p_desc = text_frame4.add_paragraph()
    p_desc.text = f"   {desc}"
    p_desc.font.size = Pt(18)
    p_desc.space_after = Pt(14)

# ========== 幻灯片 5: 核心功能 ==========
slide5 = prs.slides.add_slide(prs.slide_layouts[6])

title_box5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
title_frame5 = title_box5.text_frame
title_frame5.text = "核心功能模块"
title_para5 = title_frame5.paragraphs[0]
title_para5.font.size = Pt(44)
title_para5.font.bold = True
title_para5.font.color.rgb = SVTR_DARK

content_box5 = slide5.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(5))
text_frame5 = content_box5.text_frame
text_frame5.word_wrap = True

features = [
    "💬 智能聊天系统 - 基于RAG的AI问答",
    "📊 融资数据可视化 - 阶段、金额筛选",
    "🔍 智能建议系统 - 上下文相关推荐",
    "📈 每日交易精选 - 自动同步展示",
    "🌐 飞书知识库集成 - 252节点内容",
    "⚡ 实时数据同步 - 智能同步策略",
    "🎯 标签过滤系统 - 多维度筛选",
    "📱 响应式设计 - 完美移动体验"
]

for feature in features:
    p = text_frame5.add_paragraph()
    p.text = feature
    p.font.size = Pt(22)
    p.space_after = Pt(16)

# ========== 幻灯片 6: 开发命令 ==========
slide6 = prs.slides.add_slide(prs.slide_layouts[6])

title_box6 = slide6.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
title_frame6 = title_box6.text_frame
title_frame6.text = "开发工作流"
title_para6 = title_frame6.paragraphs[0]
title_para6.font.size = Pt(44)
title_para6.font.bold = True
title_para6.font.color.rgb = SVTR_DARK

# 左侧：开发命令
left_box6 = slide6.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(4.2), Inches(4.5))
left_frame6 = left_box6.text_frame
left_frame6.word_wrap = True

p_dev = left_frame6.paragraphs[0]
p_dev.text = "开发 & 测试"
p_dev.font.size = Pt(24)
p_dev.font.bold = True
p_dev.font.color.rgb = SVTR_BLUE

dev_cmds = [
    "npm run dev → 开发服务器",
    "npm run preview → 预览",
    "npm run test → 单元测试",
    "npm run test:e2e → E2E测试",
    "npm run lint → 代码检查"
]

for cmd in dev_cmds:
    p = left_frame6.add_paragraph()
    p.text = f"  • {cmd}"
    p.font.size = Pt(16)
    p.space_after = Pt(10)

# 右侧：构建部署
right_box6 = slide6.shapes.add_textbox(Inches(5.2), Inches(1.8), Inches(4.2), Inches(4.5))
right_frame6 = right_box6.text_frame
right_frame6.word_wrap = True

p_build = right_frame6.paragraphs[0]
p_build.text = "构建 & 部署"
p_build.font.size = Pt(24)
p_build.font.bold = True
p_build.font.color.rgb = SVTR_BLUE

build_cmds = [
    "npm run build → TS编译",
    "npm run optimize:all → 资源优化",
    "npm run sync → 飞书同步",
    "npm run deploy:cloudflare → 部署",
    "npm run backup → 备份"
]

for cmd in build_cmds:
    p = right_frame6.add_paragraph()
    p.text = f"  • {cmd}"
    p.font.size = Pt(16)
    p.space_after = Pt(10)

# ========== 幻灯片 7: 性能优化 ==========
slide7 = prs.slides.add_slide(prs.slide_layouts[6])

title_box7 = slide7.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
title_frame7 = title_box7.text_frame
title_frame7.text = "性能优化成果"
title_para7 = title_frame7.paragraphs[0]
title_para7.font.size = Pt(44)
title_para7.font.bold = True
title_para7.font.color.rgb = SVTR_DARK

# 数据展示
metrics_box = slide7.shapes.add_textbox(Inches(1.5), Inches(2), Inches(7), Inches(4))
metrics_frame = metrics_box.text_frame
metrics_frame.word_wrap = True

metrics = [
    ("📦 资源优化", "37.9KB 总体积减少"),
    ("⚡ 性能评分", "Lighthouse 90+ 分"),
    ("🖼️ 图片优化", "WebP 转换 + Fallback"),
    ("🗜️ 压缩策略", "Terser/CleanCSS + Gzip"),
    ("🚀 边缘计算", "Cloudflare 全球加速"),
    ("💾 缓存优化", "KV存储 + 智能缓存")
]

for label, value in metrics:
    p = metrics_frame.add_paragraph()
    p.text = label
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = SVTR_BLUE
    p.space_after = Pt(4)

    p_value = metrics_frame.add_paragraph()
    p_value.text = f"   {value}"
    p_value.font.size = Pt(20)
    p_value.space_after = Pt(18)

# ========== 幻灯片 8: 中文开发体验 ==========
slide8 = prs.slides.add_slide(prs.slide_layouts[6])

title_box8 = slide8.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
title_frame8 = title_box8.text_frame
title_frame8.text = "开发体验创新"
title_para8 = title_frame8.paragraphs[0]
title_para8.font.size = Pt(44)
title_para8.font.bold = True
title_para8.font.color.rgb = SVTR_DARK

content_box8 = slide8.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(5))
text_frame8 = content_box8.text_frame
text_frame8.word_wrap = True

dx_features = [
    ("🇨🇳 中文命令", "npm run 预览 | npm run 推送 | npm run 测试"),
    ("🔄 智能同步", "自动数据质量检查 + 完整性验证"),
    ("✅ 完整测试", "Jest单元测试 + Playwright E2E"),
    ("💾 自动备份", "npm run backup / rollback"),
    ("📊 质量监控", "代码质量优化 + 性能监控"),
    ("🤖 MCP集成", "多服务协同 (Feishu/GitHub/SQLite)")
]

for label, desc in dx_features:
    p = text_frame8.add_paragraph()
    p.text = label
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = SVTR_BLUE
    p.space_after = Pt(6)

    p_desc = text_frame8.add_paragraph()
    p_desc.text = f"   {desc}"
    p_desc.font.size = Pt(18)
    p_desc.space_after = Pt(16)

# ========== 幻灯片 9: 技术亮点 ==========
slide9 = prs.slides.add_slide(prs.slide_layouts[6])

title_box9 = slide9.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
title_frame9 = title_box9.text_frame
title_frame9.text = "技术亮点与创新"
title_para9 = title_frame9.paragraphs[0]
title_para9.font.size = Pt(44)
title_para9.font.bold = True
title_para9.font.color.rgb = SVTR_DARK

highlights_box = slide9.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4.5))
highlights_frame = highlights_box.text_frame
highlights_frame.word_wrap = True

highlights = [
    "🎯 混合技术栈：前端原生JS + 后端Serverless",
    "🔍 双文件系统：优化版本与源文件并存",
    "📡 实时数据流：飞书 → KV → 前端展示",
    "🧠 智能RAG：向量检索 + 语义理解",
    "⚙️ 自动化流程：测试/构建/部署一体化",
    "🔐 安全机制：API鉴权 + 数据加密"
]

for highlight in highlights:
    p = highlights_frame.add_paragraph()
    p.text = highlight
    p.font.size = Pt(22)
    p.space_after = Pt(20)

# ========== 幻灯片 10: 感谢页 ==========
slide10 = prs.slides.add_slide(prs.slide_layouts[6])

# 背景色
background10 = slide10.background
fill10 = background10.fill
fill10.solid()
fill10.fore_color.rgb = SVTR_BLUE

# 主标题
thanks_box = slide10.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
thanks_frame = thanks_box.text_frame
thanks_frame.text = "感谢观看"
thanks_para = thanks_frame.paragraphs[0]
thanks_para.font.size = Pt(64)
thanks_para.font.bold = True
thanks_para.font.color.rgb = RGBColor(255, 255, 255)
thanks_para.alignment = PP_ALIGN.CENTER

# 联系信息
contact_box = slide10.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1.5))
contact_frame = contact_box.text_frame
contact_frame.word_wrap = True

p_web = contact_frame.paragraphs[0]
p_web.text = "🌐 官网：https://svtr.ai"
p_web.font.size = Pt(24)
p_web.font.color.rgb = RGBColor(255, 255, 255)
p_web.alignment = PP_ALIGN.CENTER

p_repo = contact_frame.add_paragraph()
p_repo.text = "💻 GitHub：github.com/capmapt/chatsvtr"
p_repo.font.size = Pt(24)
p_repo.font.color.rgb = RGBColor(255, 255, 255)
p_repo.alignment = PP_ALIGN.CENTER
p_repo.space_before = Pt(12)

# 保存文件
output_path = 'ChatSVTR-功能介绍.pptx'
prs.save(output_path)
print("PPT created successfully: " + output_path)
print("Total slides: 10")
print("Location: C:/Projects/chatsvtr/" + output_path)
